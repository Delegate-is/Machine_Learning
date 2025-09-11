from django.shortcuts import render

# Create your views here.
def home(request):
    return render(request, "simplepages/home.html")

def contact_us(request):
    return render(request, "simplepages/contact_us.html")

def about_us(request):
    return render(request, "simplepages/about_us.html")

from django.shortcuts import render
from django.http import HttpResponse
from .forms import SlideForm
from pptx import Presentation
import io

def create_ppt(request):
    if request.method == "POST":
        form = SlideForm(request.POST)
        if form.is_valid():
            title = form.cleaned_data["title"]
            content = form.cleaned_data["content"]

            # Create PowerPoint
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]  # Title & Content layout
            slide = prs.slides.add_slide(slide_layout)

            slide.shapes.title.text = title
            slide.placeholders[1].text = content

            # Save to memory
            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)

            response = HttpResponse(
                pptx_io.read(),
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            response["Content-Disposition"] = 'attachment; filename="presentation.pptx"'
            return response
    else:
        form = SlideForm()

    return render(request, "simplepages/create_ppt.html", {"form": form})
